import streamlit as st
import openai
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from openai import OpenAI

st.title("事業資料自動要約AI（PoC）")

uploaded_files = st.file_uploader("PDF, PPTX, Excel をアップロード", type=["pdf", "pptx", "xlsx"], accept_multiple_files=True)

if uploaded_files:
    all_text = ""  # ← すべてのファイルからのテキストを集める

    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split('.')[-1]
        text = ""

        if file_type == "pdf":
            doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)

        elif file_type == "pptx":
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_type == "xlsx":
            df = pd.read_excel(uploaded_file, sheet_name=None)
            for name, sheet in df.items():
                text += f"【{name}】\n" + sheet.to_string() + "\n"

        all_text += text + "\n\n"  # ← すべて結合していく

    prompt = f"""
    あなたは投資家です。以下の資料テキストから次の情報を抽出・要約してください：

    1. 事業概要
    2. 売上・利益の実績と計画
    3. キャップテーブル（株主構成・調達履歴）
    4. 現在の調達希望条件とExit想定
    5. 競合他社と各社のPER
    6. Exit時のvaluation算定

    テキスト：
    {all_text[:10000]}  # 長すぎる場合の制限もここで
    """

    with st.spinner("OpenAIで要約中..."):
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        result = response.choices[0].message.content

    st.subheader("抽出結果")
    st.markdown(result)